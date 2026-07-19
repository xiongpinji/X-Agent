"""Document Processing Skills - Built-in skills for handling various document formats"""

from __future__ import annotations

import logging
from typing import Any
from pathlib import Path
import json

from .skills_core import (
    SkillMetadata,
    SkillCapability,
    SkillRiskLevel,
    SkillExecutionContext,
)

logger = logging.getLogger(__name__)


class WordDocumentSkill:
    """Skill for reading and processing Word documents (.docx)"""

    metadata = SkillMetadata(
        name="Word Document Processor",
        version="1.0.0",
        description="Read, extract text, and process Word documents (.docx)",
        author="X-Agent",
        license="MIT",
        capabilities=[SkillCapability.DOCUMENT_READ, SkillCapability.TEXT_EXTRACT],
        risk_level=SkillRiskLevel.LOW,
        timeout_seconds=60,
        max_memory_mb=256,
        tags=["document", "word", "text-extraction"],
        dependencies={"python-docx": ">=0.8.11"},
    )

    async def initialize(self) -> None:
        """Initialize the skill"""
        try:
            import docx
            self.docx = docx
        except ImportError:
            raise ImportError("python-docx is required. Install with: pip install python-docx")

    async def execute(self, context: SkillExecutionContext) -> dict[str, Any]:
        """Execute the skill"""
        file_path = context.input_data.get("file_path")
        if not file_path:
            raise ValueError("file_path is required")

        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if path.suffix.lower() != ".docx":
            raise ValueError(f"Expected .docx file, got {path.suffix}")

        # Load document
        doc = self.docx.Document(file_path)

        # Extract text
        paragraphs = [p.text for p in doc.paragraphs]
        text = "\n".join(paragraphs)

        # Extract tables
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)

        return {
            "text": text,
            "paragraphs": paragraphs,
            "tables": tables,
            "paragraph_count": len(paragraphs),
            "table_count": len(tables),
        }

    async def validate_input(self, input_data: dict[str, Any]) -> tuple[bool, str | None]:
        """Validate input data"""
        if "file_path" not in input_data:
            return False, "file_path is required"
        return True, None

    async def cleanup(self) -> None:
        """Cleanup resources"""
        pass


class ExcelSpreadsheetSkill:
    """Skill for reading and processing Excel spreadsheets (.xlsx)"""

    metadata = SkillMetadata(
        name="Excel Spreadsheet Processor",
        version="1.0.0",
        description="Read, extract data, and process Excel spreadsheets (.xlsx)",
        author="X-Agent",
        license="MIT",
        capabilities=[SkillCapability.DOCUMENT_READ, SkillCapability.DATA_ANALYZE],
        risk_level=SkillRiskLevel.LOW,
        timeout_seconds=120,
        max_memory_mb=512,
        tags=["document", "excel", "data-analysis"],
        dependencies={"openpyxl": ">=3.0.0"},
    )

    async def initialize(self) -> None:
        """Initialize the skill"""
        try:
            from openpyxl import load_workbook
            self.load_workbook = load_workbook
        except ImportError:
            raise ImportError("openpyxl is required. Install with: pip install openpyxl")

    async def execute(self, context: SkillExecutionContext) -> dict[str, Any]:
        """Execute the skill"""
        file_path = context.input_data.get("file_path")
        if not file_path:
            raise ValueError("file_path is required")

        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if path.suffix.lower() != ".xlsx":
            raise ValueError(f"Expected .xlsx file, got {path.suffix}")

        # Load workbook
        wb = self.load_workbook(file_path)

        # Extract sheet data
        sheets = {}
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_data = []
            for row in ws.iter_rows(values_only=True):
                sheet_data.append(list(row))
            sheets[sheet_name] = sheet_data

        return {
            "sheets": sheets,
            "sheet_names": wb.sheetnames,
            "sheet_count": len(wb.sheetnames),
        }

    async def validate_input(self, input_data: dict[str, Any]) -> tuple[bool, str | None]:
        """Validate input data"""
        if "file_path" not in input_data:
            return False, "file_path is required"
        return True, None

    async def cleanup(self) -> None:
        """Cleanup resources"""
        pass


class PDFDocumentSkill:
    """Skill for reading and processing PDF documents"""

    metadata = SkillMetadata(
        name="PDF Document Processor",
        version="1.0.0",
        description="Read, extract text, and process PDF documents",
        author="X-Agent",
        license="MIT",
        capabilities=[SkillCapability.DOCUMENT_READ, SkillCapability.TEXT_EXTRACT],
        risk_level=SkillRiskLevel.LOW,
        timeout_seconds=120,
        max_memory_mb=512,
        tags=["document", "pdf", "text-extraction"],
        dependencies={"PyPDF2": ">=3.0.0"},
    )

    async def initialize(self) -> None:
        """Initialize the skill"""
        try:
            from PyPDF2 import PdfReader
            self.PdfReader = PdfReader
        except ImportError:
            raise ImportError("PyPDF2 is required. Install with: pip install PyPDF2")

    async def execute(self, context: SkillExecutionContext) -> dict[str, Any]:
        """Execute the skill"""
        file_path = context.input_data.get("file_path")
        if not file_path:
            raise ValueError("file_path is required")

        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if path.suffix.lower() != ".pdf":
            raise ValueError(f"Expected .pdf file, got {path.suffix}")

        # Load PDF
        with open(file_path, "rb") as f:
            reader = self.PdfReader(f)
            pages = []
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                pages.append({
                    "page_number": page_num + 1,
                    "text": text,
                })

        # Combine all text
        all_text = "\n".join([p["text"] for p in pages])

        return {
            "text": all_text,
            "pages": pages,
            "page_count": len(pages),
        }

    async def validate_input(self, input_data: dict[str, Any]) -> tuple[bool, str | None]:
        """Validate input data"""
        if "file_path" not in input_data:
            return False, "file_path is required"
        return True, None

    async def cleanup(self) -> None:
        """Cleanup resources"""
        pass


class PowerPointSkill:
    """Skill for reading and processing PowerPoint presentations"""

    metadata = SkillMetadata(
        name="PowerPoint Presentation Processor",
        version="1.0.0",
        description="Read, extract content, and process PowerPoint presentations (.pptx)",
        author="X-Agent",
        license="MIT",
        capabilities=[SkillCapability.DOCUMENT_READ, SkillCapability.TEXT_EXTRACT],
        risk_level=SkillRiskLevel.LOW,
        timeout_seconds=120,
        max_memory_mb=512,
        tags=["document", "powerpoint", "presentation"],
        dependencies={"python-pptx": ">=0.6.21"},
    )

    async def initialize(self) -> None:
        """Initialize the skill"""
        try:
            from pptx import Presentation
            self.Presentation = Presentation
        except ImportError:
            raise ImportError("python-pptx is required. Install with: pip install python-pptx")

    async def execute(self, context: SkillExecutionContext) -> dict[str, Any]:
        """Execute the skill"""
        file_path = context.input_data.get("file_path")
        if not file_path:
            raise ValueError("file_path is required")

        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if path.suffix.lower() != ".pptx":
            raise ValueError(f"Expected .pptx file, got {path.suffix}")

        # Load presentation
        prs = self.Presentation(file_path)

        # Extract slides
        slides = []
        for slide_num, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": slide_num + 1,
                "shapes": [],
            }

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_data["shapes"].append({
                        "type": shape.shape_type,
                        "text": shape.text,
                    })

            slides.append(slide_data)

        # Combine all text
        all_text = "\n".join([
            shape["text"]
            for slide in slides
            for shape in slide["shapes"]
            if "text" in shape
        ])

        return {
            "text": all_text,
            "slides": slides,
            "slide_count": len(slides),
        }

    async def validate_input(self, input_data: dict[str, Any]) -> tuple[bool, str | None]:
        """Validate input data"""
        if "file_path" not in input_data:
            return False, "file_path is required"
        return True, None

    async def cleanup(self) -> None:
        """Cleanup resources"""
        pass


class JSONProcessorSkill:
    """Skill for processing JSON files"""

    metadata = SkillMetadata(
        name="JSON Processor",
        version="1.0.0",
        description="Read, parse, and process JSON files",
        author="X-Agent",
        license="MIT",
        capabilities=[SkillCapability.DOCUMENT_READ, SkillCapability.DATA_ANALYZE],
        risk_level=SkillRiskLevel.LOW,
        timeout_seconds=60,
        max_memory_mb=256,
        tags=["document", "json", "data-processing"],
    )

    async def initialize(self) -> None:
        """Initialize the skill"""
        pass

    async def execute(self, context: SkillExecutionContext) -> dict[str, Any]:
        """Execute the skill"""
        file_path = context.input_data.get("file_path")
        if not file_path:
            raise ValueError("file_path is required")

        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if path.suffix.lower() != ".json":
            raise ValueError(f"Expected .json file, got {path.suffix}")

        # Load JSON
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        return {
            "data": data,
            "type": type(data).__name__,
            "size": len(str(data)),
        }

    async def validate_input(self, input_data: dict[str, Any]) -> tuple[bool, str | None]:
        """Validate input data"""
        if "file_path" not in input_data:
            return False, "file_path is required"
        return True, None

    async def cleanup(self) -> None:
        """Cleanup resources"""
        pass


__all__ = [
    "WordDocumentSkill",
    "ExcelSpreadsheetSkill",
    "PDFDocumentSkill",
    "PowerPointSkill",
    "JSONProcessorSkill",
]
