"""
文档处理器 - 处理Word、PDF、PowerPoint和Excel文档
"""

import logging
from typing import Any, Dict, List, Optional
from pathlib import Path

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """文档处理器 - 支持多种文档格式"""

    def __init__(self):
        self._processors = {}
        self._initialize_processors()

    def _initialize_processors(self) -> None:
        """初始化各种文档处理器"""
        try:
            from docx import Document
            self._processors["docx"] = self._process_docx
        except ImportError:
            logger.warning("python-docx not installed, DOCX support disabled")

        try:
            import PyPDF2
            self._processors["pdf"] = self._process_pdf
        except ImportError:
            logger.warning("PyPDF2 not installed, PDF support disabled")

        try:
            from pptx import Presentation
            self._processors["pptx"] = self._process_pptx
        except ImportError:
            logger.warning("python-pptx not installed, PPTX support disabled")

        try:
            import openpyxl
            self._processors["xlsx"] = self._process_xlsx
        except ImportError:
            logger.warning("openpyxl not installed, XLSX support disabled")

    async def process(
        self,
        file_path: str,
        operation: str,
        **kwargs
    ) -> Dict[str, Any]:
        """
        处理文档

        Args:
            file_path: 文件路径
            operation: 操作类型
            **kwargs: 操作特定的参数

        Returns:
            Dict[str, Any]: 处理结果
        """
        try:
            path = Path(file_path)
            if not path.exists():
                return {"success": False, "error": f"File not found: {file_path}"}

            suffix = path.suffix.lower().lstrip(".")

            if suffix == "docx":
                return await self._process_docx(file_path, operation, **kwargs)
            elif suffix == "pdf":
                return await self._process_pdf(file_path, operation, **kwargs)
            elif suffix == "pptx":
                return await self._process_pptx(file_path, operation, **kwargs)
            elif suffix in ("xlsx", "xls"):
                return await self._process_xlsx(file_path, operation, **kwargs)
            else:
                return {"success": False, "error": f"Unsupported file format: {suffix}"}

        except Exception as e:
            logger.error(f"Error processing document: {e}")
            return {"success": False, "error": str(e)}

    async def _process_docx(
        self,
        file_path: str,
        operation: str,
        **kwargs
    ) -> Dict[str, Any]:
        """处理Word文档"""
        try:
            from docx import Document

            if operation == "extract_text":
                doc = Document(file_path)
                text = "\n".join([p.text for p in doc.paragraphs])
                return {"success": True, "data": text}

            elif operation == "add_paragraph":
                doc = Document(file_path)
                text = kwargs.get("text", "")
                doc.add_paragraph(text)
                doc.save(file_path)
                return {"success": True, "message": "Paragraph added"}

            elif operation == "add_heading":
                doc = Document(file_path)
                text = kwargs.get("text", "")
                level = kwargs.get("level", 1)
                doc.add_heading(text, level)
                doc.save(file_path)
                return {"success": True, "message": "Heading added"}

            elif operation == "get_paragraphs":
                doc = Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs]
                return {"success": True, "data": paragraphs}

            else:
                return {"success": False, "error": f"Unknown operation: {operation}"}

        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            return {"success": False, "error": str(e)}

    async def _process_pdf(
        self,
        file_path: str,
        operation: str,
        **kwargs
    ) -> Dict[str, Any]:
        """处理PDF文档"""
        try:
            import PyPDF2

            if operation == "extract_text":
                with open(file_path, "rb") as file:
                    reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text()
                return {"success": True, "data": text}

            elif operation == "get_page_count":
                with open(file_path, "rb") as file:
                    reader = PyPDF2.PdfReader(file)
                    count = len(reader.pages)
                return {"success": True, "data": count}

            elif operation == "extract_page":
                page_num = kwargs.get("page_num", 0)
                with open(file_path, "rb") as file:
                    reader = PyPDF2.PdfReader(file)
                    if page_num >= len(reader.pages):
                        return {"success": False, "error": "Page number out of range"}
                    text = reader.pages[page_num].extract_text()
                return {"success": True, "data": text}

            else:
                return {"success": False, "error": f"Unknown operation: {operation}"}

        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            return {"success": False, "error": str(e)}

    async def _process_pptx(
        self,
        file_path: str,
        operation: str,
        **kwargs
    ) -> Dict[str, Any]:
        """处理PowerPoint文档"""
        try:
            from pptx import Presentation

            if operation == "extract_text":
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return {"success": True, "data": "\n".join(text)}

            elif operation == "get_slide_count":
                prs = Presentation(file_path)
                count = len(prs.slides)
                return {"success": True, "data": count}

            elif operation == "extract_slide":
                slide_num = kwargs.get("slide_num", 0)
                prs = Presentation(file_path)
                if slide_num >= len(prs.slides):
                    return {"success": False, "error": "Slide number out of range"}
                slide = prs.slides[slide_num]
                text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                return {"success": True, "data": "\n".join(text)}

            else:
                return {"success": False, "error": f"Unknown operation: {operation}"}

        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            return {"success": False, "error": str(e)}

    async def _process_xlsx(
        self,
        file_path: str,
        operation: str,
        **kwargs
    ) -> Dict[str, Any]:
        """处理Excel文档"""
        try:
            import openpyxl

            if operation == "read_data":
                wb = openpyxl.load_workbook(file_path)
                ws = wb.active
                data = []
                for row in ws.iter_rows(values_only=True):
                    data.append(list(row))
                return {"success": True, "data": data}

            elif operation == "get_sheet_names":
                wb = openpyxl.load_workbook(file_path)
                names = wb.sheetnames
                return {"success": True, "data": names}

            elif operation == "read_sheet":
                sheet_name = kwargs.get("sheet_name")
                wb = openpyxl.load_workbook(file_path)
                if sheet_name not in wb.sheetnames:
                    return {"success": False, "error": f"Sheet not found: {sheet_name}"}
                ws = wb[sheet_name]
                data = []
                for row in ws.iter_rows(values_only=True):
                    data.append(list(row))
                return {"success": True, "data": data}

            elif operation == "write_data":
                data = kwargs.get("data", [])
                sheet_name = kwargs.get("sheet_name", "Sheet1")
                wb = openpyxl.load_workbook(file_path)
                if sheet_name not in wb.sheetnames:
                    ws = wb.create_sheet(sheet_name)
                else:
                    ws = wb[sheet_name]
                for row_idx, row in enumerate(data, 1):
                    for col_idx, value in enumerate(row, 1):
                        ws.cell(row=row_idx, column=col_idx, value=value)
                wb.save(file_path)
                return {"success": True, "message": "Data written"}

            else:
                return {"success": False, "error": f"Unknown operation: {operation}"}

        except Exception as e:
            logger.error(f"Error processing XLSX: {e}")
            return {"success": False, "error": str(e)}
